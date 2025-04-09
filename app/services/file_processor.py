import os
import logging
import tempfile
from typing import Dict, Any, Optional, Tuple, List
from fastapi import UploadFile, HTTPException, status
from datetime import datetime
import uuid
import shutil
import mimetypes
import json

from app.models.document import Document
from app.config import get_settings

# 初始化日志
logger = logging.getLogger(__name__)

# 配置获取
settings = get_settings()

# 支持的文件类型
SUPPORTED_FILE_TYPES = {
    # 文本文件
    ".txt": "text/plain",
    ".md": "text/markdown",
    ".csv": "text/csv",
    ".json": "application/json",
    
    # 文档
    ".pdf": "application/pdf",
    ".doc": "application/msword",
    ".docx": "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
    
    # 演示文稿
    ".ppt": "application/vnd.ms-powerpoint",
    ".pptx": "application/vnd.openxmlformats-officedocument.presentationml.presentation",
    
    # 表格
    ".xls": "application/vnd.ms-excel",
    ".xlsx": "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
    
    # 网页
    ".html": "text/html",
    ".htm": "text/html",
}

def get_file_extension(filename: str) -> str:
    """获取文件扩展名"""
    return os.path.splitext(filename)[1].lower()

def is_supported_file_type(filename: str) -> bool:
    """检查文件类型是否支持"""
    ext = get_file_extension(filename)
    return ext in SUPPORTED_FILE_TYPES

def guess_file_mime_type(filename: str) -> str:
    """猜测文件的MIME类型"""
    ext = get_file_extension(filename)
    if ext in SUPPORTED_FILE_TYPES:
        return SUPPORTED_FILE_TYPES[ext]
    
    # 尝试使用mimetypes模块猜测
    mime_type, _ = mimetypes.guess_type(filename)
    return mime_type or "application/octet-stream"

def save_upload_file(upload_file: UploadFile, destination_path: str) -> bool:
    """保存上传的文件到指定路径"""
    try:
        # 确保目标目录存在
        os.makedirs(os.path.dirname(destination_path), exist_ok=True)
        
        # 保存文件
        with open(destination_path, "wb") as buffer:
            shutil.copyfileobj(upload_file.file, buffer)
        
        return True
    except Exception as e:
        logger.error(f"保存文件失败: {str(e)}")
        return False
    finally:
        upload_file.file.close()

def get_file_size(file_path: str) -> int:
    """获取文件大小（字节）"""
    return os.path.getsize(file_path)

def extract_text_from_file(file_path: str, file_type: str) -> str:
    """从文件中提取文本内容"""
    try:
        # 纯文本文件处理
        if file_type in ["text/plain", "text/markdown", "text/csv"]:
            with open(file_path, "r", encoding="utf-8") as f:
                return f.read()
        
        # JSON文件处理
        elif file_type == "application/json":
            with open(file_path, "r", encoding="utf-8") as f:
                data = json.load(f)
                # 将JSON转为字符串
                return json.dumps(data, ensure_ascii=False, indent=2)
        
        # PDF文件处理
        elif file_type == "application/pdf":
            try:
                from pypdf import PdfReader
                
                reader = PdfReader(file_path)
                text = ""
                for page in reader.pages:
                    text += page.extract_text() + "\n\n"
                return text
            except ImportError:
                logger.error("未安装pypdf，无法处理PDF文件")
                return ""
        
        # Word文档处理
        elif file_type in ["application/msword", "application/vnd.openxmlformats-officedocument.wordprocessingml.document"]:
            try:
                import docx
                
                doc = docx.Document(file_path)
                text = ""
                for para in doc.paragraphs:
                    text += para.text + "\n"
                return text
            except ImportError:
                logger.error("未安装python-docx，无法处理Word文档")
                return ""
        
        # PowerPoint处理
        elif file_type in ["application/vnd.ms-powerpoint", "application/vnd.openxmlformats-officedocument.presentationml.presentation"]:
            try:
                from pptx import Presentation
                
                ppt = Presentation(file_path)
                text = ""
                for slide in ppt.slides:
                    for shape in slide.shapes:
                        if hasattr(shape, "text"):
                            text += shape.text + "\n"
                    text += "\n"
                return text
            except ImportError:
                logger.error("未安装python-pptx，无法处理PowerPoint文件")
                return ""
        
        # Excel处理
        elif file_type in ["application/vnd.ms-excel", "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet"]:
            try:
                import pandas as pd
                
                df = pd.read_excel(file_path)
                return df.to_string()
            except ImportError:
                logger.error("未安装pandas，无法处理Excel文件")
                return ""
        
        # HTML处理
        elif file_type in ["text/html"]:
            try:
                from bs4 import BeautifulSoup
                
                with open(file_path, "r", encoding="utf-8") as f:
                    soup = BeautifulSoup(f.read(), "html.parser")
                    return soup.get_text()
            except ImportError:
                logger.error("未安装beautifulsoup4，无法处理HTML文件")
                with open(file_path, "r", encoding="utf-8") as f:
                    return f.read()
        
        # 不支持的文件类型
        else:
            logger.warning(f"不支持的文件类型: {file_type}")
            return ""
    except Exception as e:
        logger.error(f"提取文本内容失败: {str(e)}")
        return ""

def process_file(
    upload_file: UploadFile,
    file_store_path: str,
    user_id: Optional[str] = None,
    metadata: Optional[Dict[str, Any]] = None
) -> Tuple[Document, str]:
    """
    处理上传的文件，提取内容，并创建文档对象
    
    参数:
        upload_file: 上传的文件
        file_store_path: 文件存储路径基础目录
        user_id: 上传用户ID
        metadata: 附加元数据
    
    返回:
        Tuple[Document, str]: 文档对象和文件实际保存路径
    """
    # 检查文件类型是否支持
    if not is_supported_file_type(upload_file.filename):
        raise HTTPException(
            status_code=status.HTTP_400_BAD_REQUEST,
            detail=f"不支持的文件类型: {get_file_extension(upload_file.filename)}"
        )
    
    # 创建文件ID和保存路径
    file_id = str(uuid.uuid4())
    safe_filename = f"{file_id}{get_file_extension(upload_file.filename)}"
    
    # 构建保存目录结构
    current_date = datetime.now().strftime("%Y/%m/%d")
    relative_path = os.path.join(current_date, safe_filename)
    file_path = os.path.join(file_store_path, relative_path)
    
    # 保存文件
    if not save_upload_file(upload_file, file_path):
        raise HTTPException(
            status_code=status.HTTP_500_INTERNAL_SERVER_ERROR,
            detail="保存文件失败"
        )
    
    # 获取文件信息
    file_size = get_file_size(file_path)
    file_type = guess_file_mime_type(upload_file.filename)
    
    # 提取文本内容
    content = extract_text_from_file(file_path, file_type)
    processed = bool(content)
    
    # 创建文档对象
    document = Document(
        id=file_id,
        filename=upload_file.filename,
        file_path=relative_path,
        file_type=file_type,
        file_size=file_size,
        content=content,
        metadata=metadata or {},
        processed=processed,
        error=None if processed else "无法提取文本内容",
        created_by=user_id,
    )
    
    return document, file_path

def batch_process_files(
    files: List[UploadFile],
    file_store_path: str,
    user_id: Optional[str] = None,
    metadata: Optional[Dict[str, Any]] = None
) -> List[Tuple[Document, str]]:
    """
    批量处理文件
    
    参数:
        files: 文件列表
        file_store_path: 文件存储基础路径
        user_id: 用户ID
        metadata: 附加元数据
    
    返回:
        List[Tuple[Document, str]]: 处理结果列表
    """
    results = []
    
    for file in files:
        try:
            result = process_file(file, file_store_path, user_id, metadata)
            results.append(result)
        except Exception as e:
            logger.error(f"处理文件 {file.filename} 失败: {str(e)}")
            # 继续处理其他文件
    
    return results 